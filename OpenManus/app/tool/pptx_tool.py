from typing import ClassVar
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import re
from app.tool.base import BaseTool


class PptxTool(BaseTool):
    name: str = "pptx_tool"
    description: str = "根据文献分析报告自动生成PPT幻灯片"
    parameters: dict = {
        "type": "object",
        "properties": {
            "report": {"type": "string", "description": "文献分析报告内容"},
            "title": {"type": "string", "description": "PPT标题"},
            "output_path": {"type": "string", "description": "输出文件路径"},
        },
        "required": ["report", "title"],
    }

    # 主题色
    THEME: ClassVar[dict] = {
        "bg": RGBColor(0x0F, 0x0F, 0x1A),        # 深夜蓝背景
        "accent": RGBColor(0x6C, 0x63, 0xFF),     # 紫色强调
        "accent2": RGBColor(0x00, 0xD4, 0xFF),    # 青色
        "text": RGBColor(0xFF, 0xFF, 0xFF),        # 白色正文
        "subtext": RGBColor(0xAA, 0xAA, 0xCC),    # 灰色副文
        "card": RGBColor(0x1E, 0x1E, 0x3A),       # 卡片背景
    }

    async def execute(self, report: str, title: str, output_path: str = None, **kwargs):
        try:
            if not output_path:
                output_path = f"uploads/{title[:20].replace(' ', '_')}.pptx"

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            # 解析报告内容
            sections = self._parse_report(report)

            # 封面
            self._add_cover_slide(prs, title)

            # 每个章节一张幻灯片
            for section in sections:
                self._add_content_slide(prs, section["title"], section["content"])

            # 结尾致谢页
            self._add_end_slide(prs)

            prs.save(output_path)
            return self.success_response(f"PPT已生成：{output_path}")

        except Exception as e:
            import traceback
            traceback.print_exc()
            return self.fail_response(f"PPT生成失败：{str(e)}")

    def _parse_report(self, report: str) -> list:
        """把 Markdown 报告按 ## 标题拆分成章节"""
        sections = []
        # 匹配 ## 或 ### 开头的标题
        parts = re.split(r'\n#{1,3} ', report)
        for part in parts[1:]:  # 跳过第一个空块
            lines = part.strip().split('\n')
            sec_title = lines[0].strip().replace('#', '').strip()
            sec_content = '\n'.join(lines[1:]).strip()
            # 清理 Markdown 符号
            sec_content = re.sub(r'\*\*(.*?)\*\*', r'\1', sec_content)
            sec_content = re.sub(r'\*(.*?)\*', r'\1', sec_content)
            if sec_title and sec_content:
                sections.append({
                    "title": sec_title[:50],
                    "content": sec_content[:800]
                })
        return sections

    def _set_bg(self, slide, color: RGBColor):
        """设置幻灯片背景色"""
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _add_textbox(self, slide, text, left, top, width, height,
                     font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT):
        from pptx.util import Inches, Pt
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color or self.THEME["text"]
        return txBox

    def _add_cover_slide(self, prs, title: str):
        slide_layout = prs.slide_layouts[6]  # 空白
        slide = prs.slides.add_slide(slide_layout)
        self._set_bg(slide, self.THEME["bg"])

        W = prs.slide_width
        H = prs.slide_height

        # 装饰色块
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.3), H
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.THEME["accent"]
        shape.line.fill.background()

        shape2 = slide.shapes.add_shape(
            1, Inches(0.4), Inches(0), Inches(0.08), H
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = self.THEME["accent2"]
        shape2.line.fill.background()

        # 标题
        self._add_textbox(
            slide, "📄 科研文献分析报告",
            Inches(1.5), Inches(2.0), Inches(10), Inches(1),
            font_size=20, color=self.THEME["accent2"], align=PP_ALIGN.LEFT
        )
        self._add_textbox(
            slide, title,
            Inches(1.5), Inches(2.8), Inches(10), Inches(2),
            font_size=32, bold=True, color=self.THEME["text"], align=PP_ALIGN.LEFT
        )
        self._add_textbox(
            slide, "AI 自动生成 · OpenManus 科研助手",
            Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
            font_size=14, color=self.THEME["subtext"], align=PP_ALIGN.LEFT
        )

    def _add_content_slide(self, prs, title: str, content: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_bg(slide, self.THEME["bg"])

        W = prs.slide_width
        H = prs.slide_height

        # 顶部标题栏
        header = slide.shapes.add_shape(
            1, Inches(0), Inches(0), W, Inches(1.2)
        )
        header.fill.solid()
        header.fill.fore_color.rgb = self.THEME["card"]
        header.line.fill.background()

        # 左侧强调条
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(0.15), Inches(1.2)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.THEME["accent"]
        bar.line.fill.background()

        # 标题文字
        self._add_textbox(
            slide, title,
            Inches(0.4), Inches(0.15), Inches(12), Inches(0.9),
            font_size=24, bold=True, color=self.THEME["text"]
        )

        # 内容卡片背景
        card = slide.shapes.add_shape(
            1, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = self.THEME["card"]
        card.line.fill.background()

        # 内容文字（按行分割，加bullet）
        lines = [l.strip() for l in content.split('\n') if l.strip()][:12]
        content_text = '\n'.join(f"• {l}" if not l.startswith('•') else l for l in lines)

        self._add_textbox(
            slide, content_text,
            Inches(0.7), Inches(1.6), Inches(12), Inches(5.4),
            font_size=16, color=self.THEME["text"]
        )

    def _add_end_slide(self, prs):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_bg(slide, self.THEME["bg"])

        W = prs.slide_width
        H = prs.slide_height

        shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.3), H)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.THEME["accent"]
        shape.line.fill.background()

        self._add_textbox(
            slide, "Thank You",
            Inches(1.5), Inches(2.5), Inches(10), Inches(1.5),
            font_size=48, bold=True, color=self.THEME["text"], align=PP_ALIGN.CENTER
        )
        self._add_textbox(
            slide, "Powered by OpenManus 科研助手",
            Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
            font_size=16, color=self.THEME["subtext"], align=PP_ALIGN.CENTER
        )